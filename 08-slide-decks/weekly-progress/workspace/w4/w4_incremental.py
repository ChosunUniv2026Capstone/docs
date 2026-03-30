import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


WORKSPACE = Path("/Users/kimhyeonseok/CodeStorage/smart-class")
ROOT = WORKSPACE / "docs" / "08-slide-decks" / "weekly-progress"
BASE = ROOT / "published" / "w3.pptx"
OUT = ROOT / "published" / "w4.pptx"

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
SKY = RGBColor(0xEA, 0xF2, 0xFB)
GRAY = RGBColor(0xEE, 0xF2, 0xF7)
LINE = RGBColor(0xC6, 0xD3, 0xE3)
TEXT = RGBColor(0x1B, 0x2A, 0x38)
MUTED = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x2E, 0x8B, 0x57)

BASE_REFS = {
    "docs": ("be55c70", "origin/main"),
    "CodexKit": ("29fd521", "HEAD"),
    "Front": ("a39d68a", "HEAD"),
    "Backend": ("e1e524b", "HEAD"),
    "PresenceService": ("eea6cc3", "HEAD"),
    "DB": ("379af2b", "HEAD"),
}


def git(repo, *args):
    return subprocess.check_output(
        ["git", "-C", str(WORKSPACE / repo), *args],
        text=True,
    ).strip()


def repo_diff_info(repo):
    base, head = BASE_REFS[repo]
    files = git(repo, "diff", "--name-only", f"{base}..{head}").splitlines()
    files = [f for f in files if f and not f.startswith("08-slide-decks/")]
    commits = git(repo, "log", "--format=%s", f"{base}..{head}").splitlines()
    shortstat = git(repo, "diff", "--shortstat", f"{base}..{head}")
    return {"repo": repo, "files": files, "commits": commits, "shortstat": shortstat}


def summarize(repo, info):
    files = info["files"]
    commits = info["commits"]
    shortstat = info["shortstat"]

    if repo == "docs":
        focus = []
        if any("07-status/implementation-roadmap.md" in f for f in files):
            focus.append("구현 로드맵 상태 문서 추가")
        if any("04-architecture/local-runtime-topology.md" in f for f in files):
            focus.append("로컬 런타임 topology 문서 추가")
        if any("05-work-items/epic-full-lms-delivery-plan.md" in f for f in files):
            focus.append("LMS delivery epic 문서화")
        if any("05-work-items/task-phase-2-academic-read-model.md" in f for f in files):
            focus.append("phase 2 academic read model task 추가")
        return {
            "title": "docs",
            "summary": " / ".join(focus[:3]) or "상태·아키텍처·작업 항목 문서 확장",
            "detail": f"{len(files)}개 파일 변경, {shortstat or '문서 범위 갱신'}",
        }

    if repo == "CodexKit":
        focus = []
        if any("install/bootstrap_workspace.sh" in f for f in files):
            focus.append("workspace bootstrap 보강")
        if any("docker-compose.yml" in f for f in files):
            focus.append("통합 로컬 runtime 구성 추가")
        if any("workspace-seed/docs/04-architecture/local-runtime-topology.md" in f for f in files):
            focus.append("docs seed에 runtime topology 반영")
        if any("validate_branch_name.py" in f for f in files):
            focus.append("짧은 브랜치 규약 전파")
        return {
            "title": "CodexKit",
            "summary": " / ".join(focus[:3]) or "운영 패키지와 seed 문서 갱신",
            "detail": f"{len(files)}개 파일 변경, 최근 커밋: {commits[-1] if commits else '없음'}",
        }

    if repo == "Front":
        return {
            "title": "Front",
            "summary": "`src/App.tsx`, `src/api.ts`, `src/index.css` 중심으로 로그인·프로필·강의 화면 재구성",
            "detail": shortstat or "프론트 콘솔 뼈대 추가",
        }

    if repo == "Backend":
        return {
            "title": "Backend",
            "summary": "`app/main.py`, `models.py`, `schemas.py`, `services.py`로 eligibility / attendance read-model spine 구축",
            "detail": shortstat or "FastAPI 백엔드 뼈대 추가",
        }

    if repo == "PresenceService":
        return {
            "title": "PresenceService",
            "summary": "`dummy_openwrt.py`, `cache.py`, `service.py`로 OpenWrt형 스냅샷 수집 더미 서비스 추가",
            "detail": shortstat or "재실성 보조 서비스 추가",
        }

    if repo == "DB":
        return {
            "title": "DB",
            "summary": "`001_schema.sql`, `010_seed.sql`, CSV seed 파일로 attendance slice 데모 데이터 구성",
            "detail": shortstat or "DB seed 추가",
        }

    raise ValueError(repo)


def set_text_frame(text_frame, text, font_size=14, bold=False, color=TEXT):
    text_frame.clear()
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.name = "Pretendard"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.line_spacing = 1.08


def add_top_title(slide, title):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.color.rgb = NAVY

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.24), Inches(0.72), Inches(0.03), Inches(0.34)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.color.rgb = NAVY

    title_box = slide.shapes.add_textbox(Inches(0.38), Inches(0.58), Inches(5.8), Inches(0.4))
    set_text_frame(title_box.text_frame, title, font_size=21, bold=True, color=TEXT)


def add_card(slide, x, y, w, h, title, body, fill_rgb=SKY):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_rgb
    box.line.color.rgb = LINE
    box.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.12), Inches(w - 0.28), Inches(0.22))
    set_text_frame(title_box.text_frame, title, font_size=13, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.42), Inches(w - 0.28), Inches(h - 0.52))
    set_text_frame(body_box.text_frame, body, font_size=10.8, color=TEXT)
    body_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def add_bullet_box(slide, x, y, w, h, title, items):
    add_card(slide, x, y, w, h, title, "", fill_rgb=GRAY)
    body_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.42), Inches(w - 0.28), Inches(h - 0.52))
    tf = body_box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Pretendard"
        p.font.size = Pt(11.2)
        p.font.color.rgb = TEXT
        p.space_after = Pt(4)
        p.line_spacing = 1.05


def add_slide_at(prs, index):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(index, new_id)
    return prs.slides[index]


def move_thanks_to_end(prs):
    thanks_idx = None
    for idx, slide in enumerate(prs.slides):
        texts = [sh.text.strip() for sh in slide.shapes if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        if "감사합니다." in texts:
            thanks_idx = idx
            break
    if thanks_idx is None:
        return
    sld_id_lst = prs.slides._sldIdLst
    thanks = sld_id_lst[thanks_idx]
    sld_id_lst.remove(thanks)
    sld_id_lst.append(thanks)


def renumber_pages(prs):
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = shape.text.strip()
            if txt.isdigit() and shape.left > 15000000 and shape.top > 9000000:
                shape.text = str(slide_index)
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.RIGHT
                    for run in paragraph.runs:
                        run.font.name = "Pretendard"
                        run.font.size = Pt(9)
                        run.font.color.rgb = MUTED


def update_cover_date(prs):
    cover = prs.slides[0]
    meta = cover.shapes[1]
    meta.text = meta.text.replace("03월 30일", "04월 06일")


def add_progress_slide(prs, docs_info, kit_info):
    slide = add_slide_at(prs, 4)
    add_top_title(slide, "팀 소개 및 진행사항 - W4 업데이트")
    add_bullet_box(
        slide,
        0.8,
        1.45,
        5.3,
        4.7,
        "문서 / 운영 변경",
        [
            f"docs: {docs_info['summary']}",
            f"CodexKit: {kit_info['summary']}",
            "기존 제안서 중심 발표를 유지하되, 이후 주차 업데이트가 가능한 운영 기반을 확보했습니다.",
        ],
    )
    add_card(slide, 6.4, 1.45, 5.95, 1.45, "docs diff", docs_info["detail"], SKY)
    add_card(slide, 6.4, 3.15, 5.95, 1.45, "CodexKit diff", kit_info["detail"], GRAY)
    add_card(
        slide,
        6.4,
        4.85,
        5.95,
        1.3,
        "배치 기준",
        "목차 01 영역에 현재 팀 진행과 운영 변화만 추가하고, 기존 후속 섹션은 유지합니다.",
        SKY,
    )


def add_repo_implementation_slide(prs, front_info, backend_info, presence_info, db_info):
    slide = add_slide_at(prs, 12)
    add_top_title(slide, "구현 진행 현황 - W4 diff 요약")
    add_card(slide, 0.75, 1.35, 3.0, 2.2, front_info["title"], f"{front_info['summary']}\n{front_info['detail']}", SKY)
    add_card(slide, 3.95, 1.35, 3.0, 2.2, backend_info["title"], f"{backend_info['summary']}\n{backend_info['detail']}", GRAY)
    add_card(slide, 7.15, 1.35, 3.0, 2.2, presence_info["title"], f"{presence_info['summary']}\n{presence_info['detail']}", SKY)
    add_card(slide, 10.35, 1.35, 2.2, 2.2, db_info["title"], f"{db_info['summary']}\n{db_info['detail']}", GRAY)
    add_bullet_box(
        slide,
        0.88,
        4.0,
        11.65,
        2.0,
        "섹션 배치 이유",
        [
            "프론트 / 백엔드 / 재실성 서비스 / DB 변화는 기존 아키텍처·구현 섹션 바로 뒤에 두는 것이 자연스럽습니다.",
            "이번 주 추가 정보는 각 저장소 diff에서 실제 생성된 파일과 기능 뼈대를 기준으로 요약했습니다.",
        ],
    )


def add_plan_slide(prs, docs_info, all_infos):
    slide = add_slide_at(prs, 18)
    add_top_title(slide, "개발 계획 보강 - W4 기준")
    add_bullet_box(
        slide,
        0.82,
        1.4,
        6.0,
        4.85,
        "현재 판단",
        [
            "Phase 1: 로그인, 단말 관리, eligibility, Docker 실행 기본 동작 확보",
            "Phase 2: 강의 목록 / 공지 / 관리자 조회 뼈대 진행",
            "후속 단계는 저장소별 첫 slice를 실제 출석 흐름으로 연결하는 작업이 중심입니다.",
        ],
    )
    repos_changed = ", ".join(info["repo"] for info in all_infos if info["files"])
    add_card(
        slide,
        7.05,
        1.4,
        5.45,
        1.45,
        "diff 기반 추적 범위",
        f"이번 W4 반영은 {repos_changed} 레포의 git diff를 기준으로 정리했습니다.",
        SKY,
    )
    add_card(
        slide,
        7.05,
        3.05,
        5.45,
        1.45,
        "docs 상태 근거",
        docs_info["summary"],
        GRAY,
    )
    add_card(
        slide,
        7.05,
        4.7,
        5.45,
        1.55,
        "다음 주 기준",
        "새 목차를 따로 늘리지 않고, 다음 업데이트도 해당 섹션 안에서 필요한 위치에만 보강합니다.",
        SKY,
    )


def main():
    infos = {repo: repo_diff_info(repo) for repo in BASE_REFS}
    summaries = {repo: summarize(repo, info) for repo, info in infos.items()}

    prs = Presentation(BASE)
    update_cover_date(prs)
    add_progress_slide(prs, summaries["docs"], summaries["CodexKit"])
    add_repo_implementation_slide(
        prs,
        summaries["Front"],
        summaries["Backend"],
        summaries["PresenceService"],
        summaries["DB"],
    )
    add_plan_slide(prs, summaries["docs"], [infos[r] for r in BASE_REFS])
    move_thanks_to_end(prs)
    renumber_pages(prs)
    prs.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()
